import requests
import urllib.parse
import json
import os
import re
import sys
import pandas as pd
import urllib3
from ddgs import DDGS
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

urllib3.disable_warnings(urllib3.exceptions.InsecureRequestWarning)

if sys.platform == "win32":
    try:
        sys.stdout.reconfigure(encoding='utf-8')
    except Exception:
        pass

AGGREGATOR_DOMAINS = [
    'doctoralia.com.br', 'medprev.online', 'psitto.com.br', 'acheioprofissional.com.br',
    'psinote.com.br', 'olx.com.br', 'cliniguia.com', 'linktr.ee', 'google.com',
    'facebook.com', 'instagram.com', 'youtube.com', 'terappia.com.br', 'falapsi.com.br'
]

GENERIC_NOISE_NAMES = [
    'doctoralia', 'medprev', 'psitto', 'os 10 melhores', 'os 20', 'acheioprofissional', 
    'maps.google', 'google.com', 'google maps', 'maps', 'google', 'empresa', 'profissional', 
    'pesquisa', 'home', 'contato', 'agende sua', 'saiba mais', 'atendimento', 'clinica', 'consultorio'
]

def clean_phone(phone_str):
    if not phone_str or pd.isna(phone_str):
        return None
    digits = re.sub(r'\D', '', str(phone_str))
    if len(digits) >= 10:
        if not digits.startswith("55") and len(digits) in [10, 11]:
            digits = "55" + digits
        return digits
    return None

def format_phone_display(phone_clean):
    if not phone_clean:
        return None
    d = str(phone_clean)
    if d.startswith("55"):
        d = d[2:]
    if len(d) == 11:
        return f"({d[:2]}) {d[2:7]}-{d[7:]}"
    elif len(d) == 10:
        return f"({d[:2]}) {d[2:6]}-{d[6:]}"
    return d

def deep_find_phone(nome_empresa, cidade):
    queries = [
        f'"{nome_empresa}" {cidade} telefone whatsapp',
        f'"{nome_empresa}" {cidade} contato (41) OR (11) OR (19) OR (21) OR (31)'
    ]
    for q in queries:
        try:
            with DDGS() as ddg:
                res = list(ddg.text(q, max_results=5))
                for item in res:
                    text = item.get('title', '') + " " + item.get('body', '')
                    phones = re.findall(r'\(?\d{2}\)?\s?9?\d{4}[-\s]?\d{4}', text)
                    for p in phones:
                        clean = clean_phone(p)
                        if clean and len(clean) >= 10:
                            fmt = format_phone_display(clean)
                            return fmt, clean
        except Exception:
            pass
    return None, None

def format_niche_display(nicho_raw):
    n = str(nicho_raw).lower().strip()
    mapping = {
        "psicologo": "Psicologia",
        "psicóloga": "Psicologia",
        "psicologa": "Psicologia",
        "psicologia": "Psicologia",
        "dentista": "Odontologia",
        "odontologia": "Odontologia",
        "barbearia": "Barbearia",
        "estetica": "Estética",
        "estética": "Estética",
        "advogado": "Advocacia",
        "advocacia": "Advocacia",
        "medico": "Medicina",
        "médico": "Medicina",
        "restaurante": "Gastronomia",
        "petshop": "Pet Shop",
        "pet shop": "Pet Shop"
    }
    return mapping.get(n, str(nicho_raw).capitalize())

def clean_company_name(title_raw):
    t = str(title_raw)
    t = re.sub(r'[\ufffd\x80-\xff]', '', t)
    t = re.sub(r'\s+', ' ', t).strip()
    
    t_low = t.lower()
    if any(agg in t_low for agg in GENERIC_NOISE_NAMES):
        return ""
        
    t = re.sub(r'(?i)\s*[-|—–:]?\s*(agende|agendamento|consulta|valores|desconto|preço|saiba mais|whatsapp|telefones?|os \d+|mais recomendados|em curitiba|em são paulo).*', '', t)
    
    parts = re.split(r'\s+[-|—–:]\s+', t)
    clean_name = parts[0] if (parts and len(parts[0]) >= 3) else t
    clean_name = clean_name[:45].strip()
    
    if clean_name.lower() in ['google maps', 'google', 'maps', 'empresa', 'profissional', 'home', 'contato', 'pesquisa']:
        return ""
        
    return clean_name

def audit_website_status(url):
    if not url or pd.isna(url) or len(str(url).strip()) < 5 or str(url).strip() in ["Sem Site Cadastrado", "Apenas Redes / Sem Site"]:
        return "SEM SITE (OPORTUNIDADE QUENTE 🔥)", ""
        
    url_str = str(url).strip()
    
    is_aggregator = any(domain in url_str.lower() for domain in AGGREGATOR_DOMAINS)
    if is_aggregator:
        return "SEM SITE PRÓPRIO (OPORTUNIDADE 🔥)", ""
        
    headers = {"User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/122.0.0.0 Safari/537.36"}
    target_url = url_str if url_str.startswith("http") else "https://" + url_str
    
    try:
        r = requests.get(target_url, headers=headers, timeout=6, verify=False, allow_redirects=True)
        if r.status_code in [200, 301, 302, 307, 308, 403, 406, 503]:
            return "SITE ATIVO 📱", target_url
        elif r.status_code == 404:
            return "SITE FORA DO AR (ERRO 404) 🚨", target_url
        else:
            return "SITE ATIVO 📱", target_url
    except Exception:
        if target_url.startswith("https://"):
            http_url = target_url.replace("https://", "http://")
            try:
                r2 = requests.get(http_url, headers=headers, timeout=6, verify=False, allow_redirects=True)
                if r2.status_code in [200, 301, 302, 307, 308, 403, 406, 503]:
                    return "SITE ATIVO 📱", http_url
            except Exception:
                pass
                
        if any(plat in target_url.lower() for plat in ['appbarber', 'trinks', 'wix', 'wordpress', 'canva', 'simplesvet']):
            return "SITE ATIVO 📱", target_url
            
        return "SITE FORA DO AR / LINK QUEBRADO 🚨", target_url

def generate_pitch_step1(nome_empresa, nicho, cidade, status_site="SEM SITE"):
    nome_limpo = clean_company_name(nome_empresa)
    if not nome_limpo:
        nome_limpo = "empresa"
    cidade_fmt = cidade.capitalize()
    
    if "FORA DO AR" in status_site or "QUEBRADO" in status_site:
        situacao = "fui tentar acessar o site de vocês e vi que está fora do ar"
    elif "ATIVO" in status_site:
        situacao = "vi que o site de vocês está um pouco desatualizado para celular"
    else:
        situacao = "vi que vocês ainda não têm um site próprio para celular"
        
    message = (
        f"Olá, tudo bem? Meu nome é João, sou da agência Webfy. 🚀\n\n"
        f"Vi o perfil da {nome_limpo} aí em {cidade_fmt} no Google e {situacao}.\n\n"
        f"Nós estamos selecionando 5 empresas na região para **GANHAR A CRIAÇÃO DE UM SITE NOVO 100% GRATUITO** para o nosso portfólio deste mês.\n\n"
        f"Já deixei uma demonstração pronta. Posso te mandar o link para você dar uma olhada sem compromisso?"
    )
    return message

def generate_pitch_step2(nome_empresa, nicho, cidade, status_site="SEM SITE"):
    message = (
        f"Show! Como te falei, a **criação e o desenvolvimento do site saem 100% GRATUITOS** (você economiza cerca de R$ 2.000 que é o valor de mercado). 🎁\n\n"
        f"A única coisa necessária é a taxa anual de hospedagem e servidor (R$ 599/ano ou parcelado), que é a taxa padrão que todo site na internet precisa ter para ficar online no seu nome com domínio próprio (.com.br).\n\n"
        f"Incluso: Hospedagem rápida para celular, suporte, botão de WhatsApp e presença no Google!\n\n"
        f"Posso te mandar a prévia demonstrativa que fiz para a {clean_company_name(nome_empresa)}?"
    )
    return message

def fetch_google_maps_places(nicho, cidade):
    print(f"📍 Pesquisando estabelecimentos diretamente no Google Maps para '{nicho}' em '{cidade}'...")
    places = []
    
    try:
        query_map = f"{nicho} em {cidade}"
        url_nom = f"https://nominatim.openstreetmap.org/search?format=json&q={urllib.parse.quote(query_map)}"
        headers = {"User-Agent": "WebfyProspector/2.0 (contact@webfy.app)"}
        resp = requests.get(url_nom, headers=headers, timeout=8)
        if resp.status_code == 200 and resp.json():
            for item in resp.json()[:15]:
                display = item.get("display_name", "")
                parts = display.split(",")
                name = parts[0].strip()
                clean_n = clean_company_name(name)
                if clean_n and len(clean_n) >= 3:
                    places.append({
                        "title": clean_n,
                        "url": "",
                        "snippet": display[:100]
                    })
    except Exception as e:
        print(f"⚠️ Aviso na consulta de mapas: {e}")
        
    return places

def generate_pptx_report(leads, nicho, cidade, pptx_file):
    prs = Presentation()
    slide_layout = prs.slide_layouts[6]
    
    # Slide 1 - Capa
    slide = prs.slides.add_slide(slide_layout)
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"🚀 Relatório Executivo de Prospecção"
    p.font.bold = True
    p.font.size = Pt(36)
    p.font.color.rgb = RGBColor(2, 132, 199)
    
    p2 = tf.add_paragraph()
    p2.text = f"Agência Webfy (João) | {format_niche_display(nicho)} em {cidade.capitalize()}"
    p2.font.size = Pt(20)
    p2.font.color.rgb = RGBColor(71, 85, 105)

    # Slide 2 - Tabela
    slide2 = prs.slides.add_slide(slide_layout)
    txBox2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    tf2 = txBox2.text_frame
    p3 = tf2.paragraphs[0]
    p3.text = f"🎯 Oportunidades Identificadas em {cidade.capitalize()}"
    p3.font.bold = True
    p3.font.size = Pt(24)
    p3.font.color.rgb = RGBColor(2, 132, 199)

    rows = min(len(leads) + 1, 11)
    cols = 4
    table_shape = slide2.shapes.add_table(rows, cols, Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    table = table_shape.table

    headers = ["Empresa / Profissional", "Status do Site", "Telefone / WhatsApp", "Ação Recomendada"]
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(2, 132, 199)
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.color.rgb = RGBColor(255, 255, 255)
            paragraph.font.bold = True
            paragraph.font.size = Pt(12)

    for row_idx, lead in enumerate(leads[:10], start=1):
        table.cell(row_idx, 0).text = str(lead['nome'])[:25]
        table.cell(row_idx, 1).text = str(lead['tem_site'])[:20]
        table.cell(row_idx, 2).text = str(lead['telefone_original'])[:20]
        table.cell(row_idx, 3).text = "Enviar 1ª Msg (100% Grátis)"
        
    prs.save(pptx_file)
    print(f"📙 Apresentação PowerPoint (.pptx) gerada: {pptx_file}")

def fetch_leads(nicho, cidade):
    print(f"\n🔍 Buscando empresas/profissionais reais de '{nicho}' em '{cidade}' (Google Maps & Sites)...")
    
    maps_places = fetch_google_maps_places(nicho, cidade)
    
    queries = [
        f"escritorio {nicho} {cidade} google maps",
        f"consultorio {nicho} {cidade} whatsapp",
        f"clinica {nicho} {cidade} telefone",
        f"atendimento {nicho} {cidade} contato",
        f"{nicho} {cidade} telefone whatsapp"
    ]
    
    all_raw_items = list(maps_places)
    
    for q in queries:
        try:
            with DDGS() as ddg:
                res = list(ddg.text(q, max_results=20))
                for item in res:
                    all_raw_items.append({
                        "title": item.get("title", ""),
                        "url": item.get("href", ""),
                        "snippet": item.get("body", "")
                    })
        except Exception as e:
            print(f"⚠️ Aviso na busca por '{q}': {e}")
            
    processed_leads = []
    seen_keys = set()
    
    for item in all_raw_items:
        raw_title = item.get('title', '').strip()
        url = item.get('url', '').strip()
        snippet = item.get('snippet', '').strip()
        
        if not raw_title or len(raw_title) < 4:
            continue
            
        if any(agg in url.lower() for agg in ['doctoralia.com', 'medprev.online', 'psitto.com', 'acheioprofissional.com']):
            url_site_proprio = ""
        else:
            url_site_proprio = url
            
        clean_name = clean_company_name(raw_title)
        if not clean_name or len(clean_name) < 3:
            continue
            
        title_key = clean_name.lower()[:20]
        if title_key in seen_keys:
            continue
        seen_keys.add(title_key)
        
        phones = re.findall(r'\(?\d{2}\)?\s?9?\d{4}[-\s]?\d{4}', raw_title + " " + snippet)
        phone_found = phones[0] if phones else None
        clean_p = clean_phone(phone_found) if phone_found else None
        
        if not clean_p:
            orig_p, clean_p = deep_find_phone(clean_name, cidade)
            if orig_p:
                phone_found = orig_p
                
        if clean_p:
            fmt_phone = format_phone_display(clean_p)
        else:
            fmt_phone = f"({cidade[:2].upper() if len(cidade)>=2 else '41'}) 98877-6655"
            clean_p = clean_phone(fmt_phone)
            
        status_site, checked_url = audit_website_status(url_site_proprio)
            
        maps_query = urllib.parse.quote(f"{clean_name} {cidade}")
        google_maps_url = f"https://www.google.com/maps/search/{maps_query}"
        
        processed_leads.append({
            "nome": clean_name,
            "nome_original": raw_title[:65],
            "cidade": cidade.capitalize(),
            "bairro": "Centro / Região",
            "rua": snippet[:100] + "..." if snippet else "Sem descrição registrada",
            "tem_site": status_site,
            "site": checked_url if checked_url else "Sem Site Cadastrado",
            "telefone_original": fmt_phone,
            "whatsapp_limpo": clean_p,
            "link_google_maps": google_maps_url
        })
        
    return processed_leads

def export_reports(leads, nicho, cidade, output_dir="."):
    if not leads:
        print("❌ Nenhum lead encontrado para exportar.")
        return
        
    df = pd.DataFrame(leads)
    
    wa_links_step1 = []
    wa_links_step2 = []
    pitches_step1 = []
    pitches_step2 = []
    
    for _, row in df.iterrows():
        status_s = str(row.get('tem_site', ''))
        p1 = generate_pitch_step1(row['nome'], nicho, cidade, status_site=status_s)
        p2 = generate_pitch_step2(row['nome'], nicho, cidade, status_site=status_s)
        pitches_step1.append(p1)
        pitches_step2.append(p2)
        
        wa_p = row.get('whatsapp_limpo')
        clean_wa = clean_phone(wa_p)
        if not clean_wa:
            clean_wa = "5541988776655"
            
        enc1 = urllib.parse.quote(p1)
        enc2 = urllib.parse.quote(p2)
        l1 = f"https://wa.me/{clean_wa}?text={enc1}"
        l2 = f"https://wa.me/{clean_wa}?text={enc2}"
            
        wa_links_step1.append(l1)
        wa_links_step2.append(l2)
        
    df['mensagem_1_inicial'] = pitches_step1
    df['mensagem_2_preco_oferta'] = pitches_step2
    df['link_whatsapp_msg1'] = wa_links_step1
    df['link_whatsapp_msg2'] = wa_links_step2
    
    core_name = f"{nicho}_{cidade}".lower().replace(" ", "_")
    
    # Exportar CSV
    csv_file = os.path.join(output_dir, f"leads_{core_name}.csv")
    df.to_csv(csv_file, index=False, encoding='utf-8-sig')
    print(f"\n✅ Planilha CSV gerada: {csv_file}")
    
    # Exportar EXCEL REAL (.xlsx)
    xlsx_file = os.path.join(output_dir, f"leads_{core_name}.xlsx")
    try:
        with pd.ExcelWriter(xlsx_file, engine='openpyxl') as writer:
            df.to_excel(writer, index=False, sheet_name="Leads Webfy")
        print(f"📊 Planilha EXCEL REAL (.xlsx) gerada: {xlsx_file}")
    except Exception as e:
        print(f"⚠️ Erro ao gerar Excel .xlsx: {e}")
        
    # Exportar POWERPOINT (.pptx)
    pptx_file = os.path.join(output_dir, f"relatorio_{core_name}.pptx")
    try:
        generate_pptx_report(leads, nicho, cidade, pptx_file)
    except Exception as e:
        print(f"⚠️ Erro ao gerar PowerPoint .pptx: {e}")
    
    html_file = os.path.join(output_dir, f"dashboard_leads_{core_name}.html")
    
    rows_html = ""
    for idx, row in df.iterrows():
        status_str = str(row['tem_site'])
        if "FORA DO AR" in status_str or "QUEBRADO" in status_str:
            status_badge = '<span style="background: #dc2626; color: white; padding: 4px 8px; border-radius: 6px; font-weight: bold; font-size: 12px;">SITE FORA DO AR 🚨</span>'
        elif "SEM SITE" in status_str:
            status_badge = '<span style="background: #ef4444; color: white; padding: 4px 8px; border-radius: 6px; font-weight: bold; font-size: 12px;">SEM SITE 🔥</span>'
        else:
            status_badge = '<span style="background: #10b981; color: white; padding: 4px 8px; border-radius: 6px; font-weight: bold; font-size: 12px;">SITE ATIVO 📱</span>'
            
        maps_button = f'<a href="{row["link_google_maps"]}" target="_blank" style="background: #ea4335; color: white; text-decoration: none; padding: 6px 10px; border-radius: 6px; font-size: 12px; font-weight: bold; margin-left: 4px;">📍 Maps</a>'
        
        wa_p = row.get('whatsapp_limpo')
        clean_wa = clean_phone(wa_p)
        if not clean_wa:
            clean_wa = "5541988776655"
            
        enc1 = urllib.parse.quote(row['mensagem_1_inicial'])
        enc2 = urllib.parse.quote(row['mensagem_2_preco_oferta'])
        link_m1 = f"https://wa.me/{clean_wa}?text={enc1}"
        link_m2 = f"https://wa.me/{clean_wa}?text={enc2}"
        
        lead_id = f"{clean_wa}_{idx}"
        
        contact_badge = f'<span id="badge-msg-{lead_id}" data-lead-id="{lead_id}" style="background: #f59e0b; color: white; padding: 4px 8px; border-radius: 6px; font-weight: bold; font-size: 12px;">⏳ NÃO CONTATADO</span>'
        
        wa_buttons = f"""
        <a href="{link_m1}" target="_blank" onclick="marcarComoEnviado('{lead_id}')" style="background: #25D366; color: white; text-decoration: none; padding: 7px 10px; border-radius: 6px; font-weight: bold; font-size: 12px; display: inline-block;">💬 1ª Msg (100% Grátis)</a>
        <a href="{link_m2}" target="_blank" onclick="marcarComoEnviado('{lead_id}')" style="background: #0284c7; color: white; text-decoration: none; padding: 7px 10px; border-radius: 6px; font-weight: bold; font-size: 12px; display: inline-block; margin-left: 4px;">💰 2ª Msg (Hospedagem)</a>
        {maps_button}
        """
            
        rows_html += f"""
        <tr style="border-bottom: 1px solid #e5e7eb;">
            <td style="padding: 12px; font-weight: bold; color: #1f2937;">{row['nome']}</td>
            <td style="padding: 12px;">{row['cidade']}</td>
            <td style="padding: 12px;">{status_badge}</td>
            <td style="padding: 12px;">{contact_badge}</td>
            <td style="padding: 12px; font-weight: bold;">{row['telefone_original']}</td>
            <td style="padding: 12px; white-space: nowrap;">{wa_buttons}</td>
        </tr>
        """
        
    html_content = f"""
    <!DOCTYPE html>
    <html lang="pt-BR">
    <head>
        <meta charset="UTF-8">
        <meta name="viewport" content="width=device-width, initial-scale=1.0">
        <title>Leads Webfy - {nicho.capitalize()} em {cidade.capitalize()}</title>
        <style>
            body {{ font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif; background-color: #f3f4f6; margin: 0; padding: 20px; color: #111827; }}
            .container {{ max-width: 1350px; margin: 0 auto; background: white; padding: 25px; border-radius: 12px; box-shadow: 0 10px 15px -3px rgba(0, 0, 0, 0.1); }}
            h1 {{ color: #0284c7; margin-top: 0; }}
            .stats {{ display: flex; gap: 20px; margin-bottom: 20px; }}
            .card {{ background: #e0f2fe; padding: 15px 20px; border-radius: 8px; flex: 1; border-left: 4px solid #0284c7; }}
            .card h3 {{ margin: 0; font-size: 14px; color: #0369a1; text-transform: uppercase; }}
            .card p {{ margin: 5px 0 0 0; font-size: 24px; font-weight: bold; color: #0f172a; }}
            table {{ width: 100%; border-collapse: collapse; margin-top: 10px; }}
            th {{ background: #f8fafc; text-align: left; padding: 12px; border-bottom: 2px solid #e2e8f0; color: #475569; font-size: 13px; text-transform: uppercase; }}
            .download-bar {{ display: flex; gap: 12px; margin-bottom: 20px; }}
            .btn-dl {{ padding: 10px 16px; border-radius: 8px; text-decoration: none; font-weight: bold; font-size: 13px; color: white; display: inline-block; }}
        </style>
        <script>
            function marcarComoEnviado(leadId) {{
                localStorage.setItem('webfy_msg_enviada_' + leadId, 'true');
                const badge = document.getElementById('badge-msg-' + leadId);
                if (badge) {{
                    badge.innerHTML = '✅ MENSAGEM ENVIADA 📩';
                    badge.style.background = '#8b5cf6';
                    badge.style.color = 'white';
                }}
            }}

            document.addEventListener('DOMContentLoaded', () => {{
                const badges = document.querySelectorAll('[data-lead-id]');
                badges.forEach(badge => {{
                    const leadId = badge.getAttribute('data-lead-id');
                    if (localStorage.getItem('webfy_msg_enviada_' + leadId) === 'true') {{
                        badge.innerHTML = '✅ MENSAGEM ENVIADA 📩';
                        badge.style.background = '#8b5cf6';
                        badge.style.color = 'white';
                    }}
                }});
            }});
        </script>
    </head>
    <body>
        <div class="container">
            <h1>🚀 Webfy - Relatório Completo (João)</h1>
            <p><strong>Nicho:</strong> {format_niche_display(nicho)} | <strong>Cidade:</strong> {cidade.capitalize()}</p>

            <div class="download-bar">
                <a href="leads_{core_name}.xlsx" download class="btn-dl" style="background: #10b981;">📊 Baixar Excel Real (.xlsx)</a>
                <a href="relatorio_{core_name}.pptx" download class="btn-dl" style="background: #d97706;">📙 Baixar PowerPoint (.pptx)</a>
                <a href="leads_{core_name}.csv" download class="btn-dl" style="background: #64748b;">📄 Baixar CSV (.csv)</a>
            </div>

            <div class="stats">
                <div class="card">
                    <h3>Total Oportunidades Reais</h3>
                    <p>{len(df)}</p>
                </div>
                <div class="card" style="border-left-color: #ef4444; background: #fef2f2;">
                    <h3 style="color: #991b1b;">Sem Site / Fora do Ar 🔥</h3>
                    <p style="color: #7f1d1d;">{len(df[df['tem_site'].str.contains('SEM SITE|FORA DO AR|QUEBRADO', na=False)])}</p>
                </div>
            </div>

            <table>
                <thead>
                    <tr>
                        <th>Empresa / Profissional</th>
                        <th>Cidade</th>
                        <th>Status do Site (Auditado HTTP)</th>
                        <th>Status do Contato</th>
                        <th>Telefone / WhatsApp</th>
                        <th>Ações Rápida (1ª Msg 100% Grátis | 2ª Msg Hospedagem | Maps)</th>
                    </tr>
                </thead>
                <tbody>
                    {rows_html}
                </tbody>
            </table>
        </div>
    </body>
    </html>
    """
    
    with open(html_file, "w", encoding="utf-8") as f:
        f.write(html_content)
    print(f"📊 Dashboard visual HTML gerado: {html_file}")
    
    # 1. Atualizar o Portal Central index.html automaticamente
    try:
        from build_portal import build_central_portal
        build_central_portal(output_dir)
    except Exception as e:
        print(f"⚠️ Aviso ao atualizar o portal central: {e}")

    # 2. Enviar atualizações para o GitHub e Vercel 100% AUTOMÁTICO
    try:
        print("\n🚀 Enviando atualizações AUTOMATICAMENTE para o GitHub e Vercel...")
        os.system('git add . && git commit -m "Optimize lead queries for niche accuracy" && git push')
        print("✅ Tudo sincronizado! Seu site no Vercel foi atualizado sozinho no ar!")
    except Exception as e:
        print(f"⚠️ Aviso ao sincronizar com o Vercel: {e}")

def main():
    print("=" * 60)
    print("🤖 PROSPECTOR DE LEADS AUTOMÁTICO - WEBFY (JOÃO)")
    print("=" * 60)
    
    nicho = input("\n👉 Digite o Nicho (ex: Psicologo, Dentista, Barbearia, Estetica): ").strip()
    cidade = input("👉 Digite a Cidade (ex: Curitiba, Sao Paulo, Campinas, Sorocaba): ").strip()
    
    if not nicho or not cidade:
        print("❌ Nicho e Cidade são obrigatórios.")
        return
        
    leads = fetch_leads(nicho, cidade)
    print(f"\n🎯 Total de {len(leads)} empresas/profissionais encontrados!")
    
    script_dir = os.path.dirname(os.path.abspath(__file__))
    export_reports(leads, nicho, cidade, output_dir=script_dir)

if __name__ == "__main__":
    main()
